#!/usr/bin/env python3
"""
PPT和PDF转Markdown工具
用于提取软件质量管理课程的PPT和PDF内容，转换为Markdown格式
"""

import os
import re
from pathlib import Path

# 检查并安装依赖
def install_dependencies():
    """安装必要的库"""
    libraries = [
        ('python-pptx', 'pptx'),
        ('pdfplumber', 'pdfplumber'),
    ]

    for package_name, import_name in libraries:
        try:
            __import__(import_name)
            print(f"✓ {package_name} 已安装")
        except ImportError:
            print(f"正在安装 {package_name}...")
            os.system(f'pip install {package_name} -q')
            print(f"✓ {package_name} 安装完成")

def extract_text_from_pptx(pptx_path):
    """从PPTX文件中提取文本内容"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation(pptx_path)
        content = []

        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            slide_title = None

            # 获取标题
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        # 判断是否是标题（根据位置或内容）
                        if not slide_title and len(text) < 100:
                            slide_title = text
                        else:
                            slide_content.append(text)

            if slide_title:
                slide_content.insert(0, f"## {slide_title}")

            if slide_content:
                content.append(f"### 第{i}页")
                content.extend(slide_content)
                content.append("")

        return "\n".join(content)
    except Exception as e:
        return f"提取PPTX失败: {str(e)}"

def extract_text_from_pdf(pdf_path):
    """从PDF文件中提取文本内容"""
    try:
        import pdfplumber

        content = []
        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text:
                    content.append(f"### 第{i}页")
                    content.append(text)
                    content.append("")

        return "\n".join(content)
    except Exception as e:
        return f"提取PDF失败: {str(e)}"

def extract_html_content(html_path):
    """从HTML课件中提取文本内容"""
    try:
        from bs4 import BeautifulSoup

        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')

        # 移除script和style标签
        for tag in soup(['script', 'style']):
            tag.decompose()

        # 获取文本
        text = soup.get_text(separator='\n', strip=True)

        # 清理空行
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        return '\n'.join(lines)
    except Exception as e:
        return f"提取HTML失败: {str(e)}"

def main():
    """主函数"""
    base_dir = Path(r"e:\Codes\2026\2026-SQM")
    output_dir = base_dir / "extracted_content"
    output_dir.mkdir(exist_ok=True)

    print("=" * 60)
    print("PPT/PDF 转 Markdown 工具")
    print("=" * 60)

    # 安装依赖
    print("\n[1/3] 检查并安装依赖...")
    install_dependencies()

    # 尝试安装BeautifulSoup
    try:
        from bs4 import BeautifulSoup
        print("✓ BeautifulSoup 已安装")
    except ImportError:
        print("正在安装 BeautifulSoup...")
        os.system('pip install beautifulsoup4 -q')
        print("✓ BeautifulSoup 安装完成")

    print("\n[2/3] 提取内容...")

    # 处理PPT文件
    ppt_dir = base_dir / "课件"
    if ppt_dir.exists():
        print("\n--- 处理PPT文件 ---")
        for pptx_file in ppt_dir.glob("*.pptx"):
            print(f"正在处理: {pptx_file.name}")
            content = extract_text_from_pptx(pptx_file)
            output_file = output_dir / f"{pptx_file.stem}.md"
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# {pptx_file.name}\n\n")
                f.write(content)
            print(f"  ✓ 已保存: {output_file.name}")

    # 处理PDF文件
    pdf_dirs = [
        base_dir / "往年真题",
        base_dir / "复习背诵"
    ]

    for pdf_dir in pdf_dirs:
        if pdf_dir.exists():
            print(f"\n--- 处理目录: {pdf_dir.name} ---")
            for pdf_file in pdf_dir.glob("*.pdf"):
                print(f"正在处理: {pdf_file.name}")
                content = extract_text_from_pdf(pdf_file)
                output_file = output_dir / f"{pdf_file.stem}.md"
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(f"# {pdf_file.name}\n\n")
                    f.write(content)
                print(f"  ✓ 已保存: {output_file.name}")

    # 处理HTML课件
    html_dirs = [
        base_dir / "课件" / "XP",
        base_dir / "课件" / "课堂测试"
    ]

    for html_dir in html_dirs:
        if html_dir.exists():
            print(f"\n--- 处理目录: {html_dir.name} ---")
            for html_file in html_dir.glob("*.html"):
                print(f"正在处理: {html_file.name}")
                content = extract_html_content(html_file)
                output_file = output_dir / f"{html_file.stem}.md"
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(f"# {html_file.name}\n\n")
                    f.write(content)
                print(f"  ✓ 已保存: {output_file.name}")

    print("\n[3/3] 完成!")
    print(f"\n所有提取的内容已保存到: {output_dir}")
    print("=" * 60)

if __name__ == "__main__":
    main()
